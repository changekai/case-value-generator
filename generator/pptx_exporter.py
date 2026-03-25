# -*- coding: utf-8 -*-
"""
PPTX 簡報匯出模組
使用 python-pptx 生成 9 張投影片的簡報
"""

from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from config import PPTX_COLORS, PPTX_FONTS, PPTX_FONT_SIZES, ORG_NAME, PROJECT_NAME


# ─── 工具函式 ──────────────────────────────────────────

def _rgb(hex_str: str) -> RGBColor:
    """將 Hex 色碼轉為 RGBColor"""
    return RGBColor.from_string(hex_str)


def _set_font(run, size_pt: int, bold: bool = False, color_hex: str = None, font_name: str = None):
    """設定文字格式"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = font_name or PPTX_FONTS["body"]
    if color_hex:
        run.font.color.rgb = _rgb(color_hex)


def _add_textbox(slide, left, top, width, height, text: str,
                 size: int = 16, bold: bool = False, color: str = None,
                 alignment=PP_ALIGN.LEFT, font_name: str = None):
    """在投影片上新增文字框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold, color, font_name)
    return txBox


def _add_filled_rect(slide, left, top, width, height, fill_hex: str):
    """在投影片上新增填色矩形"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.fill.background()
    return shape


def _add_footer(slide, page_num: int):
    """在投影片左下角加入頁碼與組織名稱"""
    footer_text = f"P.{page_num}　｜　{ORG_NAME}"
    _add_textbox(
        slide,
        Inches(0.5), Inches(6.8), Inches(5), Inches(0.4),
        footer_text,
        size=PPTX_FONT_SIZES["footer"],
        color=PPTX_COLORS["gray"],
    )


def _add_paragraph(text_frame, text: str, size: int = 16, bold: bool = False,
                    color: str = None, space_before: Pt = None, alignment=PP_ALIGN.LEFT):
    """在文字框架中新增段落"""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    run = p.add_run()
    run.text = text
    _set_font(run, size, bold, color)
    return p


# ─── 投影片生成函式 ───────────────────────────────────

def _slide_1_cover(prs, meta):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面

    # 深藍背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(PPTX_COLORS["primary"])

    # 標題
    _add_textbox(
        slide, Inches(1), Inches(2), Inches(8), Inches(1),
        "潛在案源價值說明",
        size=PPTX_FONT_SIZES["cover_title"], bold=True,
        color=PPTX_COLORS["white"],
        alignment=PP_ALIGN.CENTER,
    )

    # 副標題
    subtitle = f"{meta['company_name']} — {meta['industry']}｜{meta['transformation_topic']}"
    _add_textbox(
        slide, Inches(1), Inches(3.2), Inches(8), Inches(0.8),
        subtitle,
        size=PPTX_FONT_SIZES["subtitle"], bold=False,
        color=PPTX_COLORS["secondary"],
        alignment=PP_ALIGN.CENTER,
    )

    # 組織名稱
    _add_textbox(
        slide, Inches(1), Inches(4.5), Inches(8), Inches(0.6),
        f"{PROJECT_NAME}",
        size=PPTX_FONT_SIZES["body"], bold=False,
        color=PPTX_COLORS["secondary"],
        alignment=PP_ALIGN.CENTER,
    )

    # 日期（右下角）
    date_str = datetime.now().strftime("%Y-%m-%d")
    _add_textbox(
        slide, Inches(7), Inches(6.5), Inches(2.5), Inches(0.4),
        date_str,
        size=PPTX_FONT_SIZES["small"],
        color=PPTX_COLORS["gray"],
        alignment=PP_ALIGN.RIGHT,
    )


def _slide_2_methodology(prs):
    """Slide 2: 研究方法說明 — 2M2C + 4T 框架"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 標題
    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "研究方法：2M2C + 4T 分析框架",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    # 四個框架區塊
    items = [
        ("M1", "全球趨勢洞察", "Macro Environment", "Trend\n科技/經濟/法令/應用趨勢"),
        ("M2", "台灣在地市場", "Market Environment", "Topic\n選題方向"),
        ("C1", "社會需求與價值", "Customer", "Trade\n社會需求程度"),
        ("C2", "潛在業者生態", "Competitor", "Target\n目標業者清單"),
    ]

    for i, (code, name, eng, detail) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 4.8)
        top = Inches(1.5 + row * 2.5)

        # 背景框
        rect = _add_filled_rect(slide, left, top, Inches(4.3), Inches(2), PPTX_COLORS["secondary"])

        # 代號
        _add_textbox(
            slide, left + Inches(0.2), top + Inches(0.15), Inches(1), Inches(0.5),
            code, size=22, bold=True, color=PPTX_COLORS["accent"],
        )
        # 名稱
        _add_textbox(
            slide, left + Inches(0.8), top + Inches(0.15), Inches(3), Inches(0.5),
            name, size=18, bold=True, color=PPTX_COLORS["primary"],
        )
        # 英文 + 對應 4T
        _add_textbox(
            slide, left + Inches(0.3), top + Inches(0.7), Inches(3.7), Inches(1.2),
            f"{eng}\n\n{detail}", size=13, color=PPTX_COLORS["dark_text"],
        )

    # 底部說明
    _add_textbox(
        slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
        "三個敘事維度：What（為什麼選這個題目）→ Who（為什麼是這家業者）→ We（資策會能提供什麼）",
        size=PPTX_FONT_SIZES["small"], bold=False, color=PPTX_COLORS["gray"],
        alignment=PP_ALIGN.CENTER,
    )

    _add_footer(slide, 2)


def _slide_3_m1(prs, m1):
    """Slide 3: M1 全球趨勢（四格並列）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "WHAT｜M1 全球趨勢洞察",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    # 四格趨勢
    trends = [
        ("🔬 科技趨勢", m1.get("technology_trend", "")),
        ("💰 經濟趨勢", m1.get("economics_trend", "")),
        ("📋 法令政策", m1.get("politics_trend", "")),
        ("📱 應用趨勢", m1.get("application_trend", "")),
    ]

    for i, (label, content) in enumerate(trends):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 4.8)
        top = Inches(1.3 + row * 2.2)

        # 標籤
        _add_textbox(
            slide, left, top, Inches(4.3), Inches(0.4),
            label, size=14, bold=True, color=PPTX_COLORS["accent"],
        )
        # 內容
        rect = _add_filled_rect(slide, left, top + Inches(0.4), Inches(4.3), Inches(1.5), PPTX_COLORS["secondary"])
        _add_textbox(
            slide, left + Inches(0.15), top + Inches(0.5), Inches(4), Inches(1.3),
            content[:200] if content else "（待補充）",
            size=12, color=PPTX_COLORS["dark_text"],
        )

    # 綜合論述
    summary_text = m1.get("summary", "")
    if summary_text:
        _add_textbox(
            slide, Inches(0.5), Inches(5.9), Inches(9), Inches(0.8),
            f"📌 {summary_text[:300]}",
            size=11, bold=False, color=PPTX_COLORS["primary"],
        )

    _add_footer(slide, 3)


def _slide_4_m2(prs, m2):
    """Slide 4: M2 台灣市場（左文右圖式）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "WHAT｜M2 台灣在地市場分析",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    # 左側：市場規模 + PEST
    _add_textbox(
        slide, Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.4),
        "📊 市場規模", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    _add_textbox(
        slide, Inches(0.5), Inches(1.7), Inches(4.3), Inches(1),
        m2.get("market_size", "（待補充）")[:200],
        size=12, color=PPTX_COLORS["dark_text"],
    )

    _add_textbox(
        slide, Inches(0.5), Inches(2.7), Inches(4.3), Inches(0.4),
        "🔍 PEST 分析摘要", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    _add_textbox(
        slide, Inches(0.5), Inches(3.1), Inches(4.3), Inches(1.2),
        m2.get("pest_summary", "（待補充）")[:250],
        size=12, color=PPTX_COLORS["dark_text"],
    )

    # 右側：供需缺口
    right_x = Inches(5.3)

    _add_textbox(
        slide, right_x, Inches(1.3), Inches(4.3), Inches(0.4),
        "⬆ 供給端缺口", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    supply_text = "\n".join([f"• {g}" for g in m2.get("supply_gap", ["待補充"])])
    rect = _add_filled_rect(slide, right_x, Inches(1.7), Inches(4.2), Inches(1.3), PPTX_COLORS["secondary"])
    _add_textbox(
        slide, right_x + Inches(0.15), Inches(1.8), Inches(3.9), Inches(1.1),
        supply_text, size=12, color=PPTX_COLORS["dark_text"],
    )

    _add_textbox(
        slide, right_x, Inches(3.1), Inches(4.3), Inches(0.4),
        "⬇ 需求端缺口", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    demand_text = "\n".join([f"• {g}" for g in m2.get("demand_gap", ["待補充"])])
    rect2 = _add_filled_rect(slide, right_x, Inches(3.5), Inches(4.2), Inches(1.3), PPTX_COLORS["secondary"])
    _add_textbox(
        slide, right_x + Inches(0.15), Inches(3.6), Inches(3.9), Inches(1.1),
        demand_text, size=12, color=PPTX_COLORS["dark_text"],
    )

    # 底部選題方向
    _add_textbox(
        slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.4),
        "📌 選題方向", size=14, bold=True, color=PPTX_COLORS["primary"],
    )
    _add_textbox(
        slide, Inches(0.5), Inches(5.6), Inches(9), Inches(1),
        m2.get("topic_rationale", "（待補充）")[:300],
        size=12, color=PPTX_COLORS["dark_text"],
    )

    _add_footer(slide, 4)


def _slide_5_c1c2(prs, c1, c2):
    """Slide 5: C1 社會需求 + C2 業者生態（上下兩段）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "WHAT｜C1 社會需求 + C2 業者生態",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    # 上半部：C1
    _add_textbox(
        slide, Inches(0.5), Inches(1.2), Inches(2), Inches(0.4),
        f"社會需求程度：{c1.get('social_need_score', '中')}",
        size=14, bold=True, color=PPTX_COLORS["accent"],
    )

    reasons_text = "\n".join([f"• {r}" for r in c1.get("social_need_reasons", [])])
    _add_textbox(
        slide, Inches(0.5), Inches(1.7), Inches(4.3), Inches(1.3),
        reasons_text, size=12, color=PPTX_COLORS["dark_text"],
    )

    _add_textbox(
        slide, Inches(5.3), Inches(1.2), Inches(4.3), Inches(0.4),
        "社會價值", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    _add_textbox(
        slide, Inches(5.3), Inches(1.7), Inches(4.3), Inches(1.3),
        c1.get("social_value", "（待補充）")[:250],
        size=12, color=PPTX_COLORS["dark_text"],
    )

    # 分隔線
    _add_filled_rect(slide, Inches(0.5), Inches(3.3), Inches(9), Inches(0.03), PPTX_COLORS["gray"])

    # 下半部：C2
    _add_textbox(
        slide, Inches(0.5), Inches(3.5), Inches(4.3), Inches(0.4),
        f"主營運商：{c2.get('primary_operator_type', '')}",
        size=14, bold=True, color=PPTX_COLORS["accent"],
    )

    partners_text = "\n".join(
        [f"• {p.get('type', '')}: {p.get('role', '')}" for p in c2.get("partner_types", [])]
    )
    _add_textbox(
        slide, Inches(0.5), Inches(4), Inches(4.3), Inches(1.5),
        partners_text, size=12, color=PPTX_COLORS["dark_text"],
    )

    _add_textbox(
        slide, Inches(5.3), Inches(3.5), Inches(4.3), Inches(0.4),
        f"潛在輔導規模：{c2.get('target_scale', '')}",
        size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    _add_textbox(
        slide, Inches(5.3), Inches(4), Inches(4.3), Inches(1.5),
        c2.get("competitive_landscape", "（待補充）")[:250],
        size=12, color=PPTX_COLORS["dark_text"],
    )

    _add_footer(slide, 5)


def _slide_6_company(prs, meta):
    """Slide 6: 業者介紹"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 上半部深藍底
    _add_filled_rect(slide, Inches(0), Inches(0), Inches(10), Inches(3.5), PPTX_COLORS["primary"])

    _add_textbox(
        slide, Inches(1), Inches(1), Inches(8), Inches(0.8),
        "WHO：為什麼是這家業者？",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["white"],
        alignment=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide, Inches(1), Inches(2), Inches(8), Inches(0.8),
        meta["company_name"],
        size=36, bold=True,
        color=PPTX_COLORS["white"],
        alignment=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide, Inches(1), Inches(3), Inches(8), Inches(0.5),
        f"{meta['industry']}　｜　{meta['transformation_topic']}",
        size=PPTX_FONT_SIZES["subtitle"],
        color=PPTX_COLORS["secondary"],
        alignment=PP_ALIGN.CENTER,
    )

    # 補充說明
    if meta.get("supplement"):
        _add_textbox(
            slide, Inches(1), Inches(4.2), Inches(8), Inches(2),
            meta["supplement"][:400],
            size=14, color=PPTX_COLORS["dark_text"],
        )

    _add_footer(slide, 6)


def _slide_7_feasibility(prs, who):
    """Slide 7: 五項可行性評估（五格卡片）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "WHO｜五項可行性評估",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    items = [
        ("📈 市場可行性", who.get("market_feasibility", "")),
        ("💼 商業模式", who.get("business_model", "")),
        ("👥 目標客戶", who.get("target_customer", "")),
        ("🛠 服務可行性", who.get("service_feasibility", "")),
        ("🤖 解決方案可行性", who.get("solution_feasibility", "")),
    ]

    # 第一行3個，第二行2個
    positions = [
        (0.3, 1.3), (3.5, 1.3), (6.7, 1.3),
        (1.9, 4.2), (5.1, 4.2),
    ]

    for i, ((label, content), (x, y)) in enumerate(zip(items, positions)):
        _add_filled_rect(slide, Inches(x), Inches(y), Inches(2.8), Inches(2.5), PPTX_COLORS["secondary"])
        _add_textbox(
            slide, Inches(x + 0.15), Inches(y + 0.1), Inches(2.5), Inches(0.4),
            label, size=12, bold=True, color=PPTX_COLORS["accent"],
        )
        _add_textbox(
            slide, Inches(x + 0.15), Inches(y + 0.5), Inches(2.5), Inches(1.8),
            content[:150] if content else "（待補充）",
            size=10, color=PPTX_COLORS["dark_text"],
        )

    _add_footer(slide, 7)


def _slide_8_value(prs, who):
    """Slide 8: 題目價值與策略方向"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "WHO｜題目價值與策略方向",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    # 價值主張（重點標注）
    _add_filled_rect(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(2.5), PPTX_COLORS["secondary"])
    _add_textbox(
        slide, Inches(0.7), Inches(1.3), Inches(2), Inches(0.4),
        "💎 題目整體價值主張", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    _add_textbox(
        slide, Inches(0.7), Inches(1.8), Inches(8.5), Inches(1.7),
        who.get("value_proposition", "（待補充）")[:400],
        size=13, color=PPTX_COLORS["dark_text"],
    )

    # 外溢效益
    _add_textbox(
        slide, Inches(0.5), Inches(4), Inches(2), Inches(0.4),
        "🌊 外溢效益", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    _add_textbox(
        slide, Inches(0.5), Inches(4.4), Inches(9), Inches(1),
        who.get("spillover_benefit", "（待補充）")[:250],
        size=12, color=PPTX_COLORS["dark_text"],
    )

    # 策略方向
    _add_textbox(
        slide, Inches(0.5), Inches(5.5), Inches(2), Inches(0.4),
        "🎯 策略方向", size=14, bold=True, color=PPTX_COLORS["accent"],
    )
    directions = who.get("strategic_directions", [])
    for i, d in enumerate(directions[:3]):
        _add_filled_rect(
            slide,
            Inches(0.5 + i * 3.1), Inches(5.9),
            Inches(2.8), Inches(0.8),
            PPTX_COLORS["primary"],
        )
        _add_textbox(
            slide,
            Inches(0.6 + i * 3.1), Inches(6),
            Inches(2.6), Inches(0.6),
            d[:50], size=11, bold=True, color=PPTX_COLORS["white"],
            alignment=PP_ALIGN.CENTER,
        )

    _add_footer(slide, 8)


def _slide_9_we(prs, we):
    """Slide 9: 資策會輔導價值（流程時間軸）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.7),
        "WE｜資策會能提供什麼？",
        size=PPTX_FONT_SIZES["slide_title"], bold=True,
        color=PPTX_COLORS["primary"],
    )

    _add_textbox(
        slide, Inches(0.5), Inches(1), Inches(9), Inches(0.5),
        we.get("headline", ""),
        size=PPTX_FONT_SIZES["subtitle"], bold=True,
        color=PPTX_COLORS["accent"],
        alignment=PP_ALIGN.CENTER,
    )

    # 時間軸 — 5 個階段橫向排列
    stages = we.get("stages", [])
    stage_width = 1.7
    gap = 0.15
    start_x = 0.3

    for i, stage in enumerate(stages):
        x = start_x + i * (stage_width + gap)

        # 圓形編號
        _add_filled_rect(
            slide,
            Inches(x + 0.5), Inches(1.8),
            Inches(0.6), Inches(0.6),
            PPTX_COLORS["accent"],
        )
        _add_textbox(
            slide, Inches(x + 0.5), Inches(1.85),
            Inches(0.6), Inches(0.5),
            str(i + 1), size=18, bold=True,
            color=PPTX_COLORS["white"],
            alignment=PP_ALIGN.CENTER,
        )

        # 連接線（除了最後一個）
        if i < len(stages) - 1:
            _add_filled_rect(
                slide,
                Inches(x + stage_width + 0.1), Inches(2.05),
                Inches(gap + 0.1), Inches(0.05),
                PPTX_COLORS["gray"],
            )

        # 階段名稱
        _add_textbox(
            slide, Inches(x), Inches(2.6),
            Inches(stage_width), Inches(0.5),
            stage.get("stage", ""),
            size=12, bold=True, color=PPTX_COLORS["primary"],
            alignment=PP_ALIGN.CENTER,
        )

        # 說明
        _add_filled_rect(
            slide, Inches(x), Inches(3.2),
            Inches(stage_width), Inches(2.5),
            PPTX_COLORS["secondary"],
        )
        _add_textbox(
            slide, Inches(x + 0.1), Inches(3.3),
            Inches(stage_width - 0.2), Inches(1.2),
            stage.get("description", "")[:80],
            size=10, color=PPTX_COLORS["dark_text"],
        )

        # 產出物
        outputs = "、".join(stage.get("outputs", []))
        _add_textbox(
            slide, Inches(x + 0.1), Inches(4.6),
            Inches(stage_width - 0.2), Inches(1),
            f"📦 {outputs}",
            size=9, color=PPTX_COLORS["gray"],
        )

    _add_footer(slide, 9)


# ─── 主要匯出函式 ──────────────────────────────────────

def export_pptx(analysis: dict, output_path: str):
    """
    將分析結果匯出為 PPTX 簡報檔案。

    Args:
        analysis: build_analysis() 的回傳結果
        output_path: 輸出檔案路徑

    Returns:
        輸出檔案路徑
    """
    prs = Presentation()

    # 設定投影片尺寸為 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    meta = analysis["meta"]
    m1 = analysis["m1_macro"]
    m2 = analysis["m2_market"]
    c1 = analysis["c1_customer"]
    c2 = analysis["c2_competitor"]
    who = analysis["who_evaluation"]
    we = analysis["we_contribution"]

    # 生成 9 張投影片
    _slide_1_cover(prs, meta)
    _slide_2_methodology(prs)
    _slide_3_m1(prs, m1)
    _slide_4_m2(prs, m2)
    _slide_5_c1c2(prs, c1, c2)
    _slide_6_company(prs, meta)
    _slide_7_feasibility(prs, who)
    _slide_8_value(prs, who)
    _slide_9_we(prs, we)

    prs.save(output_path)
    return output_path
